
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(17,20,24)
BG2 = RGBColor(28,33,39)
TEXT = RGBColor(247,241,232)
MUTED = RGBColor(182,188,196)
ORANGE = RGBColor(255,122,26)
TEAL = RGBColor(25,197,183)
GOLD = RGBColor(255,209,102)
RED = RGBColor(255,94,91)

slides = [
    ("Bird-safe wind energy needs trustworthy AI", "Not just a model: a credible way of building deployable AI for a difficult real-world problem."),
    ("Why this problem matters", "Wind energy and biodiversity share the same airspace."),
    ("What made our team different", "We built a disciplined research workflow, not a one-off leaderboard model."),
    ("System architecture", "Movement + radar + context into one robust decision pipeline."),
    ("What actually moved the needle", "Ranking-aware modeling, shift-aware validation, and judgment over tricks."),
    ("Why we trust this result", "The workflow increased the quality and credibility of the final system."),
    ("How this gets used in the real world", "Decision support for monitoring and targeted mitigation.")
]

for i,(title,subtitle) in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = BG
    for top,color in [(0,ORANGE),(7.32,TEAL)]:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(13.333), Inches(0.18))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(8.8), Inches(1.4))
    tf = tx.text_frame; p=tf.paragraphs[0]; run=p.add_run(); run.text=title; run.font.size=Pt(28); run.font.bold=True; run.font.color.rgb=TEXT; run.font.name='Segoe UI Semibold'
    tx2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(8.8), Inches(0.8))
    p2 = tx2.text_frame.paragraphs[0]; r2=p2.add_run(); r2.text=subtitle; r2.font.size=Pt(16); r2.font.color.rgb=MUTED; r2.font.name='Segoe UI'
    tag = slide.shapes.add_textbox(Inches(10.9), Inches(0.25), Inches(1.8), Inches(0.3))
    tp = tag.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.RIGHT; tr=tp.add_run(); tr.text='AI Cup 2026 finalist pitch'; tr.font.size=Pt(9); tr.font.color.rgb=MUTED; tr.font.name='Segoe UI'

# slide-specific simple visuals
s=prs.slides[0]
for (l,w,h,c) in [(8.7,2.55,2.55,BG2),(10.0,1.55,1.55,ORANGE),(8.2,1.25,1.25,TEAL)]:
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(1.15 + (10.0-l)*0.4), Inches(w), Inches(h)); shp.fill.solid(); shp.fill.fore_color.rgb=c; shp.line.fill.background()
for (l,text,c) in [(0.7,'Biodiversity',ORANGE),(2.05,'Wind energy',TEAL),(3.55,'Deployable AI',BG2)]:
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(2.95), Inches(1.25 if text!='Deployable AI' else 1.55), Inches(0.33)); shp.fill.solid(); shp.fill.fore_color.rgb=c; shp.line.fill.background()
    tb=s.shapes.add_textbox(Inches(l), Inches(3.0), Inches(1.55), Inches(0.2)); p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=text; r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=TEXT; r.font.name='Segoe UI Semibold'

# slide 2 three cards
s=prs.slides[1]
card_data=[('Ecology','Bird strikes are a real conservation issue.',ORANGE),('Operations','Mitigation has to be targeted and defensible.',TEAL),('Monitoring','Radar is scalable, but classification is noisy and seasonally shifting.',GOLD)]
for idx,(t,b,c) in enumerate(card_data):
    left=0.7+idx*3.9
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(3.3), Inches(2.2)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=c
    tb=s.shapes.add_textbox(Inches(left+0.18), Inches(2.18), Inches(2.8), Inches(0.4)); p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=TEXT
    tb2=s.shapes.add_textbox(Inches(left+0.18), Inches(2.7), Inches(2.9), Inches(1.1)); p2=tb2.text_frame.paragraphs[0]; r2=p2.add_run(); r2.text=b; r2.font.size=Pt(12); r2.font.color.rgb=MUTED

# slide 3 workflow boxes
s=prs.slides[2]
items=['Hypothesis','Experiment','Validation','Keep / reject','Deployable system']
subs=['Start from an idea','Run controlled tests','Check month-shift robustness','Preserve what generalizes','Turn insights into pipeline']
for i,(t,b) in enumerate(zip(items,subs)):
    left=0.55+i*2.48
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.3), Inches(2.0), Inches(1.8)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=ORANGE if i%2==0 else TEAL
    tb=s.shapes.add_textbox(Inches(left+0.1), Inches(2.55), Inches(1.8), Inches(0.35)); p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=TEXT
    tb2=s.shapes.add_textbox(Inches(left+0.1), Inches(3.0), Inches(1.8), Inches(0.7)); p2=tb2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run(); r2.text=b; r2.font.size=Pt(10.5); r2.font.color.rgb=MUTED

# slide 4 pipeline
s=prs.slides[3]
flow=[('Inputs','Trajectories + context',ORANGE),('Feature stack','Kinematics + signatures + catch22',TEAL),('Models','Ranking-aware ensemble',GOLD),('Output','Probabilistic decision support',RED)]
for i,(t,b,c) in enumerate(flow):
    left=0.8+i*2.55
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.5), Inches(2.0), Inches(1.85)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=c
    tb=s.shapes.add_textbox(Inches(left+0.1), Inches(2.8), Inches(1.8), Inches(0.3)); p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=TEXT
    tb2=s.shapes.add_textbox(Inches(left+0.1), Inches(3.15), Inches(1.8), Inches(0.7)); p2=tb2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run(); r2.text=b; r2.font.size=Pt(10.5); r2.font.color.rgb=MUTED
call=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(2.25), Inches(1.65), Inches(2.0)); call.fill.solid(); call.fill.fore_color.rgb=BG2; call.line.color.rgb=TEAL
ct=s.shapes.add_textbox(Inches(11.15), Inches(2.45), Inches(1.35), Inches(1.5)); p=ct.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text='Critical constraint\n\nSeasonal / month shift'; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=TEXT

# slide 5 insight cards
s=prs.slides[4]
ins=[('Ranking-aware', 'Metric alignment changed modeling choices.', ORANGE), ('Shift-aware', 'Local gains often failed across months.', TEAL), ('Judgment', 'We rejected brittle methods.', GOLD)]
for idx,(t,b,c) in enumerate(ins):
    left=0.7+idx*3.9
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.05), Inches(3.25), Inches(2.2)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=c
    tb=s.shapes.add_textbox(Inches(left+0.18), Inches(2.25), Inches(2.8), Inches(0.4)); p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text=t; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=TEXT
    tb2=s.shapes.add_textbox(Inches(left+0.18), Inches(2.8), Inches(2.8), Inches(1.0)); p2=tb2.text_frame.paragraphs[0]; r2=p2.add_run(); r2.text=b; r2.font.size=Pt(12); r2.font.color.rgb=MUTED

# slide 6 proof blocks
s=prs.slides[5]
proof=[('Finalist-level','Reached the AI Cup 2026 finals.',ORANGE),('Experiment scale','Broad exploration, not one lucky run.',TEAL),('Most submissions','Disciplined throughput and learning.',GOLD),('Honest evaluation','Cross-month thinking improved credibility.',RED)]
for idx,(t,b,c) in enumerate(proof):
    left=0.55+idx*3.15
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.15), Inches(2.65), Inches(2.15)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=c
    tb=s.shapes.add_textbox(Inches(left+0.15), Inches(2.35), Inches(2.3), Inches(0.5)); p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=TEXT
    tb2=s.shapes.add_textbox(Inches(left+0.18), Inches(2.95), Inches(2.25), Inches(0.8)); p2=tb2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run(); r2.text=b; r2.font.size=Pt(11); r2.font.color.rgb=MUTED

# slide 7 deployment flow
s=prs.slides[6]
flow2=[('Radar monitoring',ORANGE),('AI ranking + uncertainty',TEAL),('Operator / ecologist review',GOLD),('Targeted mitigation action',RED)]
for idx,(t,c) in enumerate(flow2):
    left=0.65+idx*3.15
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.45), Inches(2.5), Inches(1.3)); shp.fill.solid(); shp.fill.fore_color.rgb=BG2; shp.line.color.rgb=c
    tb=s.shapes.add_textbox(Inches(left+0.12), Inches(2.83), Inches(2.2), Inches(0.3)); p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=TEXT

out = Path(r'G:\Projects\epoch-ai-cup\docs\presentation\ai-cup-congress-pitch-v2.pptx')
prs.save(out)
print(f'Saved: {out}')
